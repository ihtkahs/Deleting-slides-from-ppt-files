import os
from pptx import Presentation

def modify_and_save_presentations(input_folder, output_folder):
    for ppt_file in os.listdir(input_folder):
        if ppt_file.endswith(".pptx") or ppt_file.endswith(".ppt"):
            ppt_path = os.path.join(input_folder, ppt_file)
            output_path = os.path.join(output_folder, ppt_file)

            # Load the PowerPoint presentation
            presentation = Presentation(ppt_path)

            # Delete all slides except the last 3
            slides_to_delete = len(presentation.slides) - 3
            presentation.slides._sldIdLst[:slides_to_delete] = []

            # Save the modified presentation to the output folder
            presentation.save(output_path)
            print(f"Presentation modified and saved to {output_path}")

# Replace 'input_folder' and 'output_folder' with your folder paths
modify_and_save_presentations('Jan 26 - Feb 1', 'Jan 26 - Feb 1/Student copy')
